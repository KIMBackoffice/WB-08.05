# src/export_pptx.py
"""
Export schedule as PowerPoint — hospital digital-signage style.

Structure (one .pptx per month):
    Deckblatt  ->  Abschnitt Woche 1 (Datum-Datum)  ->  Karten Woche 1 (3/Slide)
               ->  Abschnitt Woche 2                ->  Karten Woche 2  ->  ...

Card colour by Zielgruppe (HAL_Monitore Pausenraeume):
    Gelb     F8CF3B  Aerzteschaft
    Rosa     EE867A  Pflege / NDS
    Gruen    009870  Aerzteschaft + Pflege / NDS
    Hellblau 6EC1E4  PAs (Pflegeassistenten)   [Platzhalter-Farbe, anpassbar]

Text per card:
    Zeile 1  Wochentag, DD. Monat YYYY, HH:MM-HH:MM Uhr   (Arial 18, fett)
    Zeile 2  Thema (Verantwortlicher)                     (Arial 14)
    Zeile 3  Ort, Zielgruppe: ...                          (Arial 14)

Dependencies: python-pptx  (das Karten-Icon ist als base64 eingebettet).
"""
from __future__ import annotations

import base64
import os
import tempfile

import pandas as pd

from src.zielgruppe import EVENT_ZIELGRUPPE
from src.constants import MONTH_NAMES_DE

# ── colours ──────────────────────────────────────────────────────────────────
_YELLOW = "F8CF3B"   # Aerzteschaft
_SALMON = "EE867A"   # Pflege / NDS
_GREEN  = "009870"   # Aerzteschaft + Pflege / NDS
_BLUE   = "6EC1E4"   # PAs — Platzhalter, frei anpassbar
_GREY   = "EDEDED"   # Folienhintergrund
_WHITE  = "FFFFFF"
_BLACK  = "1A1A1A"

_WHITE_TEXT_CARDS: set = set()        # Referenz: schwarzer Text auf ALLEN Karten (auch grün)


def _text_color(bg_hex: str) -> str:
    return _WHITE if bg_hex in _WHITE_TEXT_CARDS else _BLACK


WEEKDAY_DE_FULL = {
    "Monday": "Montag", "Tuesday": "Dienstag", "Wednesday": "Mittwoch",
    "Thursday": "Donnerstag", "Friday": "Freitag",
    "Saturday": "Samstag", "Sunday": "Sonntag",
}

# Zielgruppen-Beschriftung fuer Zeile 3 (Reihenfolge A, P, S, PA)
_AUDIENCE_LABELS = [
    ("A",  "Ärzteschaft"),
    ("P",  "Pflegende ICU"),
    ("S",  "Studierende NDS"),
    ("PA", "Pflegeassistenten"),
]

# Sonderveranstaltungen mit fixer Zielgruppen-Beschriftung
_SPECIAL_AUDIENCE = {
    "KimSim":           "für angemeldete MAs",
    "Sitzungen_Pflege": "GL / SL / BL Pflege",
}

# ── embedded card icon (weisse Kachel mit Notiz + Stift) ─────────────────────
_ICON_B64 = "iVBORw0KGgoAAAANSUhEUgAAAgAAAAIACAYAAAD0eNT6AABZlElEQVR4nO3dd3gUdf4H8PfM7O7spveEbsGCSrGA0uTO8lPRsyuKcAhS7F1OWEiWkKjoWU8PExQL9sZ5FhS8Ext2RFFUOBFCSc+mbXZnd2fm98cmmwxSAimzm32/nsfnOdL2u7kHPu/5ls8XICIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIi6lqC2QPoiX744Qfd7DEQEfU0Q4YMYc3qRPxldgALPRGR+RgMDgx/afuBBZ+IKPIxELQPf0n7wKJPRBS9GAb2jL+Y3eho0R88eHBnDYWIiJqtX7++Q9/PMGDEX0Yb+1v4WeiJiMy3v8GAQSAk5n8J+1P0WfCJiCLf/gSCWA4DMfvG21v4WfSJiKJXe8NALAaBmHvDQPuKPws/EVHP0Z4gEGshIKbe7L4KP4s+EVHPt68wECtBICbeJAs/ERHtKtaDQI9+c8Deiz8LPxER7S0I9OQQ0GPfGLDn4s/CT0REu9pTEOipIaBHvik+9RMR0YGIpdmAHvVmAD71ExFRx8XCbIBo9gA6E4s/ERF1hj3VjZ50P0yPCQAs/kRE1Jl6egjoEVMZu/s/g4WfiIg6y+6WBKJ9OSCqB8+nfiIi6i49bV9A1C4BsPgTEVF36mlLAlEbAHaHxZ+IiLpST6ozURkAuOZPRERm2V29icZZgKgLACz+RERktp4QAqIqALD4ExFRpIj2EBA1AYDFn4iIIk00h4CoCAAs/kREFKmiNQRERQDYFYs/ERFFkmisSxEfAKIhRREREe0q0utXRAcATv0TEVG0iLalgIgNACz+REQUbaIpBERsANgViz8REUWDaKlXERkAIjUtERERHYhIrGsRGQB2FS1pioiICIiOuhVxAWDXlBQNv0QiIqJd7Vq/Im0WIKICQKT9coiIiDpTJNW5iAoAu+LTPxERRbNIrmMREwAiKRURERF1lUipdxETAHYVyamJiIiovSK1nkVsACAiIqKuExEBgDv/iYioJ4vEEwEREQCIiIioe5keACIhBREREXU3s+uf6QFgV5z+JyKinijS6lvEBQAiIiLqeqYGAG7+IyKiWBJJmwE5A0BERBSDGACIiIhikGkBgNP/REQUiyJlGYAzAERERDGIAYCIiCgGMQAQERHFIFMCANf/iYgolkXCPgDOABAREcUgBgAiIqIYxABAREQUgxgAiIiIYhADABERUQzq9gDAEwBERETmnwTgDAAREVEMYgAgIiKKQQwAREREMYgBgIiIKAYxABAREcUgBgAiIqIYxABAREQUgxgAiIiIYhADABERUQxiACAiIopBDABEREQxiAGAiIgoBjEAEBERxSAGACIiohjEAEBERBSDGACIiIhiEAMAERFRDGIAICIiikEMAERERDGIAYCIiCgGMQAQERHFIAYAIiKiGMQAQEREFIMYAIiIiGIQAwAREVEMYgAgIiKKQQwAREREMYgBgIiIKAYxABAREcUgBgAiIqIYxABAREQUgxgAiIiIYpDF7AFQdBk0aIjZQyDqND///IPZQyAyDWcAiChmMdBSLOMMAO0T/5GknmzQoCGcCaCYxBkAIop5DLkUixgAiIjAEECxhwGAiKgZQwDFEu4BoAPGdVOKNu0p8NwTQLGCMwBEFPPef3+F4c+cCaBYwABARASGAIo9DABERM0YAiiWMAAQEbXBEECxggGAiGgXDAEUCxgAiIh2gyGAejoGACKiPWAIoJ6MAYCIaC8YAqinYgAgItoHhgDqiRgAiIjagSGAehoGACKidmIIoJ6EAYCIaD8wBFBPwQBARLSfGAKoJ2AAICI6AAwBFO0YAIiIDhBDAEUzBgAiog5gCKBoxQBARNRBDAEUjRgAiIg6AUMARRsGACKiTsIQQNGEAYCIqBMxBFC0YAAgIupkDAEUDRgAiIi6AEMARToGACKiLsIQQJGMAYCIqAsxBFCkYgAgIupiDAEUiRgAiIi6AUMARRoGACKibsIQQJGEAYCIqBsxBFCkYAAgIupmDAEUCRgAiIhMwBBAZmMAICIyCUMAmYkBgIjIRAwBZBYGACIikzEEkBkYAIiIIgBDAHU3BgAiogjBEEDdiQGAiCiCMARQd2EAICKKMAwB1B0YAIiIIhBDAHU1BgAiogjFEEBdiQGAiCiCMQRQV2EAICKKcAwB1BUYAIiIogBDwJ5pmmb2EKISAwARUZSI9RAQDAZ3+/GWAMAgsH8sZg+AiIja7/33V+CMM84K/3nQoCH4+ecfTBxR1wsGg7BYLGhs9CAlJRllZeXYtm07fvnlF/zvf7/ht99+Q27uPCiKgoMPPggJCQlmDzkqMAAQEUWZWAsBbnctMjMzMGHCFSgpKYHd7oCmafD7FVgsFgiCgPPPvxhWqw05OdlYtKgQw4YNNXvYEY9LAEREUSiWlgPOPfdCHH30UOzcWQpRFOHzeREI+CEIAoLBIAKBAHRdRzDox44dO3D55ZPx6quvmz3siMcAQEQUpWIhBMydOx8eTxMAAcFgILzOr+s6dF0HAAjNX6tpOlRVhSAIyM+/C4WFi8wZdJTgEgARxby20+nRrictBzz99DK8/fYKBAL+8MdUTUdA02C3SKjxKqjy+nBkegrKGpuQHe8AgObZgABefvlV/O9/m1BQkI8+fXqb9TYiFmcAiIh6mJ4wE/DZZ5/jvvvuNxR/XzAISRRwUHICTsjJwLXHDcLT55yMx88cjdMP7oNt9Y2GnxEI+PHNN9/h0kuvwOrVH3f3W4h4nAEgIuqBonkmoLKyEvPm5aF1cj/0VF/r82PJ+BNx8XGDAF1HfW0D4qwWWGQrBmem4q41CXj8u5+RFeeAVQo93waDAdTWunHddTfh5ptvwIwZ00x6V5GHMwBERBRRnM5cVFVVQ9PU8Me2NXhw8/BjcPExA1FTXYuamjoENA0N/gDcDU2wSiLuPf0k3HvKCAQ1DXVK68yBpmnQNBX/+Mc/MXv2nfD7A2a8rYjDAEBEMSNan4hjyb333o8vvvgawWBrkS5rbMLFRxyMO0cNRb3HC0EQIAgCRKF1hsCvanB7vLh6xGAsGT8GR6QlY2djk+FnBwJ+vPfeB5gwYSJ++mlDt72nSMUAQEQxhSEgci1f/iaWLXvBsO5fr/hxdGYq5owaCkkQoDbv/N8Td4MHpw/sjyfGj8XFRxyEkvrG8GkBIBQCNm36DZdeOhH/+te/u+y9RAPuASCimNPTQkBP2PT3008bMG9eLtrW96CmQRAEzB01DEdmpcHt8bbrZ7k9XvRPikfRWWNwSEoiHvjqRyTKVjgsoZKnqkEIgoD58xdg06b/4Y47bu2KtxTxOANARESm8vv9cDpzIQiS4Wm9rNGLO0cOwfgjD2538W/RGAjCEwjC+acT8OgZo5Bos6KqyRf+fMtRwWXLXsDMmdegoqKy095PtGAAICIiU82dm4vNm7dAVVsv+9lW78GMYUfg+uHHwN3gOaCfq+k63A1NuHzo4Xhi/FiM6puN7fXGnxUI+PH5519hwoSJ+Oyzzzv0PqINAwAREZmmuPhJrFz5gWHdv6rJh9MO6o25o4ahSen4jn13QxNO7JOFJ8aPwfRhR2B7vQfBNjcHBoMBVFRUYebMa/DMM8s6/HrRggGAiIhM8eGHH+Hhhx81FH9vIIheCXGYO2ooMuLsUFR1Lz+h/dxeBfFWKx78v5G4+08nwBtQUW84KqhC0zTcf//DzT0Iej4GACIi6nY7duzAvHm5AFrX/HVdR0MggLmjhmLEgF5w+5ROfU1FVeFu8uH6kUNRPH40Dk5JRNlujgr++9/vYMKEK7Bp0/869fUjDQMAERF1O6czD/X1jeHLfYBQs5/bRgzGZUMPh7uhaS/f3THuBg/GH34Qnhg/Fn85rP9ujwpu2PAzLr10It59970uG4fZGACIiKhb3XjjbVi7dp2h2c/OxiZcftQh+NtJQ1DXuH87/g+E2+PFYalJWHLWGNxx4hCUe7xQgq3LDcFgEIrix+zZc/Hww492+XjMwABARETdQlFCU/pud7Vh3b/O58dx2emYM3IYNF2Hto9mP52l3h+Aompw/Xk4Hj59JGySiBpv67KDrmtQ1SCefPJpXH/9zaitreuWcXUXBgAiIuoWvuY1fU3TIYrNl/VoGqySiHmjh+HQjBQ0BoJ7+xGdTtV1uBub8NfjBuGJ8WNxfK8MbG/441HBjz/+BJdeOhFff/1tt46vKzEAEBFRl/N4PBg37lSMG3cq1q//Kbz2b5NEDMtOw6nHDER1Y9et+++Lu8GDMf174YnxYzBl8GHYVt8IVWu7LyCAnTt3YsqUaXjhhZdMG2dnYgAgIqIu53TmQdN01NTUGqb/A6qODVW1gE8xXO5jBrfXhzS7jEfPGIWFJx+Per8fjW1uDlSbjyTec899WLCgwKxhdhoGACIi6lKPPfY4/vvf1QgE/IaNfwAgCsDvtQ3YWdtgaM5jFm9QRa1XwS2jj0XRmaPROyEO5W3aEOu6jkAggNdf/xcmTZqCLVu2mjjajmEAICKiLrNy5QdYvLjI8NTflkUU0TshDkpQhUUUDcfxzKIjtCRw3lGH4onxY3HGIX12e1Tw++9/xKWXXoFVq/5j3mA7gAGAiIi6xJYtW5pv+NtzUW/Z9b+t3oPUpHgIJi8DtOVubMLRmal44qyxuHn4MSj1eOFX2x4VDMDjacQtt9yOf/6zyMSRHhgGACIi6hJOZx68Xp+h2Y8kSYav0QEENB33f/UjPv1tO1IT47t5lHtXp/ihAyg8dQQeOPVECBAMHQo1TYOqqnj88SW45Zbb4PGYt5FxfzEAEBFRp3O58rF+/U8IBluP9VVVViIxMRE7d+wwhII0h4wvdlRg+ruf4Nm1PyM1IQ6WCJoJCGga3I1eXHXC0VgyfgwGZ6Zhx26OCv7nPx9hwoSJWLfue5NGun8YAIiIqFM9//xLeOONNw3r/vX19Rg8dCgeWbwY1910E+rr6uDz+cKfT3PI8Ksablr1OfI+/ApWSUSSzWrG8PfI3eDBnw/ugyfGj8HEow5FSX2joWlRIODHli1bcfnlk/Hqq6+bONL2YQAgIqJO8/XX36Cw8G4E2jT0CQQCiIuLw9Tp05GaloaZV1+NOfPnIyk5Ge6amvDXyRYJ2fEO/P3L9Zix4lNsrKlDarzDjLexR+4mH3olxGHxmaORN+ZYuH0Kmtq8V1VVIQgC8vPvQmHhIhNHum8MAERE1Clqa2vhdOY2d/kLPRnruo7a2lpcNWMGxo4bh+qqKpSXl+Oc885Dbn4+jh8+HGVlZeGfIQgC+icl4O1NJZj+7qd459ctEbcvwBMIok7xY/bY4/HPM0Yj3SGjssl4VDAYDODll1/F1KnTsWPHThNHu2cMAERE1CmczlyUlZWHG+YAQHlZGS6bOBGXXH45ytsU+sqKChx11FHIzc/HJRMmoKK83LBfICchDlvrGjBzxaf4x+ffIzXODvsuGwjN1HJU8OJjBuKJ8WPxp/69sK2+0fA1gYAf33zzHS699HKsXv2xOQPdCwYAIiLqsIce+gc++WQNAoHWRj/V1dUYd8opmDZjBmrd7j98T319PaxWK2bPmYObbrsNwUAAjY2tRTRRtiHOasHc1d/glpWfozEQQKpD7pb3017uxiYcl5OOJ8aPxbXHHYUdDR4E1NYNjsFgALW1dbjhhlvwzDPLTBzpHzEAEBFRh7zzzgo8+eRThk1/TU1NGDBgAK6aMQOyLBue7tsKBAKorq7GX6dOxfwFC9B/wABUVlaGP28RRfRNiscT637F9Hc/wZc7KpCaGNfl72l/uH1+WCUR955+Eu49ZQQCmoY6pfV3EToqGDQsdUQCBgAiIjpgmzZtgtOZC7XNU2/L2fhpM2fiqGOOMTzV70l5WRlO/tOfkJufj1NPOw1lpaWGBkJ9k+KxZnsFpr/zCV5ctxGpiXGm3x3Qll/V4PZ4cfWIwVgyfgyOSEvGzgZjT4Dnnnsxok4HMAAQEdEB0XUdTmceVFWFrrcGgMqKClw5bRrOOvtsVLV5mt+X6upq9O7dG/NcLkydMQPumhooSmvTnYw4OxoCAVy/cg0KVn+DeKsFCTZLp76njnI3ePB/A/vjifFj8dfBA1HfPBOg6zpUVcXvv28xd4BtMAAQEdEBcTpz8csvGw3T+2WlpTj3ggsw+corDZv+2qupqQk+nw833Hwz/uZ0wuFwoLa2Nvx5h8WCzDg77v78e8xc8Sm21jVG3lFBjxeHpifj7IH9UONtDTCSJGHz5t9MHJkRAwAREe23p556Fm+/vcKw7h8MBDDsuOMwe84ceL3evXz33mmahoryclxw8cXIzc/HkKFDDWGi5ajgG79uwfR3P8XK/5VE3FFBrxLAoSlJ6JMUD6l5qUJVVWzcGAoAe9oT0Z0YAIiIaL/88MN6PPLIY3+44U+yWNDQ0IB1a9ciJTW1w69TUV6Ooccei9z8fJx/0UUoLyszHDHslRCHjTV1mPHup3j8y/VIjXfAJplf1nRdR2MggDirBdVNPujNPRFEUYTHE2ohvOudCGYw/zdFRERRZciQwRg3bixsNpvh44IgoLy0FHfefjuWLlmC1LS0P3zN/qqrrUVcXBzudDpx/c03w+f1hosoACTLNlhFEbM//AqzV32BgKoh1d6x1+woQRAgCgKSZRtEQYC/eYOkqqrhDZGRcOshAwAREe23hx76O+66ayEAYzFLSEyE3eHAow89hEWFhWhqakJySkqHXsvv98PtdmPajBmY53KhV+/eqK6qCn/eKonokxiPf67dgOnvfoK1ZdVITTD3qKAoCEjITMXYfjnhJQAAkGU7Nm7cZLgMySyRtX2SqJ3mzFlg9hDIZHffnWf2EGLSW2+9gxUr3sPGjZuQnp4Oi8WKYDBg+BpJkpCdk4Plr72G7du24aqZMzH8xBNRUV7eodcuLyvDqaefjj79+mFpcTH++8EHyM7JCQeQfkkJWF1Sis21DZgzaiguGXw46ho8MKPUpibG4/HVX2Pl7zswIDkBQGgJ4JxzzsLhhx9mwoj+iAGAiIj2admy5/Hhh6uxbt16BAJ+aJqG0tKyvT7JZufk4Id165Cfm4tpM2bggosvRlVlZYeefquqqjDgoIMwf8EC9OnbFy8+/zySkpLCSw2ZcQ5UexVc+94abHY3YPZJQxDQNHgC3bfpLjXegVUbt+DuNd+jt2EmQsC2bdsAhJYDzN4HwABARER79dNPG3DPPfdCEMTmhj/6Pr+nRUpqKpqamrCosBA7tm/HlVddBSB03O9AeRobIUkSbr79dvTt1w9PFhejrq4OycnJAIA4qwV2i4T8z77D5toGzB01FANSk+Bu8u3jJ3dcgtWC36pqcdea7xHQNMRZW8usKIo4+eSxALgJkIiIIpzf74fTmQtBkJp34DfvaBcAiwSomo5AUIcgAIp/98FAlmWkpqVh6ZIlKHC5sHPnTqSnp3doXKqqorKiApdedhly8/Mx6KijDEsMYvNRwRc3/Ibp736KD3/f0eVHBUVBgCAIuGvNOqwtr0ay3LoZ0Wq14bLLLsXUqVO6dAz7gwGAiIj2aO7cPGzevAWq2jqFXlGjwePTUdeoo7RKR0qigBMGSUhJFFBRs/vpfUEQkNOrF/7zwQfIz83Fxx99hOycnA6Pr7y8HCeMGIHc/Hycc+65KCszLkv0SYzH+soaTH/3Ezz5zU9ITXDAKnZN6UtOcGDR59/jpZ83G6b+LRYrjj9+GJzOv3XJ6x4oLgEQEdFuFRc/iZUrVxnO+1e6NZx+khX33uRAerIASQTsAyTACvy0OoC8Ii9WfBZA/5zdT3FnZmaiZMsWLMzNxfZt2zBx0iTU19cbbhHcX7VuN5JTUjBn/nz06dcPzyxdCqvFAkdcqAin2mX4VRW3fvAlNtfWY87IYUiWbYYLezoqNTEOL36/EQ989SP6tZlpEEURSUkJKCjI77TX6iycASAioj/48MOP8PDDjxpv+PPp6J0pYu40O/ofIUEUAFUDlBIN3k0qjj7GgucWxuOGy+zYVq4hENz9kkBCYiIsVisevv9+3HfPPQgEAkhKSurQeBWfD/X19Zh5zTVw5uYiPSMDNdXV4c/bJAm9Ehx46OufMOPdT/BTpbvTjgqmOmR8ubUUd61Zh0TZussZfwGFhQvRp0/vTnmtzsQAQEREBjt27MC8eblou9lP13U0NulwzXJgyAgLlEodoijAIoWKnSgKUOp02G0C7pkbh4duc0AJAHWNuw8BFosFWdnZeOWll5Cfm4ufN2xAZlZWh8at6zrKy8pw5vjxyF24ECeNGmW4VbClhfD7v+/AVe9+gn9t+A2pifHoSEseWZJQ6fHirjXfo6zRC4eldWLdarXh5ptvwJ/+dHKH3ldXYQAgIiIDpzMP9fWNhrX0kjINd15px4UX2aCU7/kUgBIAlGod069yYFl+PA7vL2Jn5Z6P/eXk5ODbr79Gfm4u3n7zTWRnZ3e4S15lZSUOO+ww5Obn4/JJk1BVWWlYYsiOd6C0sQlXv/cZHvjsO6Q4ZDgsB7YrP0624q413+O/W3YiI84e/rjVasOZZ56GGTOmdei9dCXuAaCoxCYwRF2jsPBurF27ztDcZ0eFhsnjZdwxw4FAZfuOACplGv50ihXP9RPhKvLihff86J8j7ra4p6aloa62FncvXIjt27fjymnTEAwGO3ShUENDAywWC27729/Qt18/LF2yBA319UhsXmpIsFmhajpyP/42dFRw5FD0TkqA29v+o4KpifH4x+ff44l1v6JvUuu6vyRZcOihB6OgYOEBj787cAaAiIgAAK+88hpeeuk1w7p/bYOG4UdLcM2yA5qO/enho1Tp6NtXxBP58cid4UB5jQ6vsvsAYXc4kJScjCWLF6NwwQJUVlYiNS2tQ+8nGAyiqrISEydPxvwFC3DoYYehsqIi/HlJFNAvKQHPrN+E6e9+ik9LStt9VDA13oF3f96MRZ//gJyE1uuIBUGArmsoKFgAm83aofF3NQYAIiLCunXfIy8v33DcLxDUIdsEuGY60OsQCcoB9O5RPEDQB9x5qwPF8+IQ7xBQVbv7FCGKInJ69cJ7K1YgPzcXX6xZ0zlHBcvKMGr0aOTm5+OMs84K7Qtok2T6Jsbj27IqTH/3Ezy79mekJsTBspdliCSbFb9U1OCuz7+HDh2WNscKJcmCwsJ8HH30UR0ed1djACAiinEejwdOZy4kSQpvmAOA8modeTPtGHuaDUp1+7v/7UrVAKVcxyWXyHhuYTxGDbWgpGzPUwlZWVn4bdMm5Ofm4oVly5CRmQmLpWMr1jU1NcjMzIQzLw8zrrkGdXV18Plap/vTHDL8qoabVn0O14dfwSqJSLJZoes6dF1HUNOg6TqsooigpuGuNevwU6UbSbs0+/nrX6/A+eef26GxdhcGACKiGOd05mHbth3Nnf5CtpaquPFyGVMn2aHspVjvD6Vcx7HHhY4KXn2RjJIyDaq6+2DRslb/wL334oF77w19LDGxQ6/v9XrR2NiIa6+/HnPmz0dSUhLcNTXhz8sWCdnxDtz35XrMXPEpNrrrkZaZirTEeNgtEtKTE5EgW/HINxvw+q9bkLNLs5+RI0/EHXfc2qExdiduAiQiimGPPfY4/vvf1YZ1//IaDX852QbXLAf0ugN/8t8dpVZHYoKA++fHYWA/EQuf8EEQdCTF/3HK3Wq1IiMzEy8sWxa+VfCYwYNRWVl5wK+v6zrKy8txznnnoU/fvli6ZAnWfPYZcpqXGlqOCr61qQRb6hpx7mH98Ut1HTbX1qPRH0SV14dGf2CXZj8SMjLSUVAQXbeUcgaAiChGrVz5ARYvLjIUf49Xx8G9Rbhm2WFNEOA/8AZ9e6QogOLWcc0sB55ZEIeDeokordp7C+Ev1qxBfm4u3nv3XcMVwAeqsqICRx19NHLz83HJhAmoKC9HMNi6/yEnIQ6baxuwZN2v+NfGLfi9thFNgSBkSUKaXd7l9XUUFuYjMzOjQ2PqbgwAREQxaMuWLZg3L9ew5q9poV36C2Y5MOhYC5T6zn3635VSpuH0M2x4riAe542zYmupahhPW2np6aisrERhfj6KFy9GUnIyZLt9t1/bXvX19bBarZg9Zw5uuu02BAIBNDY2hj+faLNCFARkxjlgt0hQ2zQUamG12jB79m0YNeqkDo3FDAwAREQxyOnMg9frMzT72VauwTnNjnPOl6FUdG3xb6FU6jj4YAnPLIzHnVMdKK3S93irYFxcHBISErD40Udx98KFqKutRUpqaodePxAIoLq6Gn+dOhXzFyxA/wEDUNWOJQaLxYqcnBwMH34CpkyZ3KExmIV7AIiIYozLlY/1638yTHlvL9cw7TwZN1/lgL+8czb9tZfSqMMiAbl3ODCwrwhXsRfVdRrSk//4jCqKInJycvD2m2+G9wWMGj0a5W2uAj4Q5WVlGPenP4WaBhUXY+V77+1xqUGSLLjhhmswc+b0Dr2m2TgDQEQUQ55//iW88cabhnX/mjoNo4dZ4JrpAAI69jAL36WCKqBU6Jh4hR3PLYzHCYMs2LaXIJKVnY2fN2xAfm4uXnnpJWRmZUGSDqydb4vq6mr07t0b81wuTJ0+HVVVVX+4pTAUCHQEAqHw1DZERRsGACKiGPHVV1+jsPDucPECACWgIzFegGuWHRn9RCgH3n23UyhlGkacaMVzBfGYdq4NW0tVaNruE0lycjICgQDuu+cePHz//bBYrYhPSOjQ6zc1NcHn82Gey4Xrb7wRVquxm5+u6xAEARkZ6QDQ4f4EZmIAICKKAW63G05nHkRRRMstf7quo8oduuHvpHE2KDUmPPrvhuLWkZYq4JG8eNxzQxzcDaGbCHfHZrMhPT0dzz71FBbm5qJkyxZkZHRsN74oivj999/x008/oa621vA5q9WG6dOnYcKESzr0GpGAAYCIKAbMnZuH8vJyQ7OfkjINt06SMelyudOa/XQWxQf463XceJ0DT7vi0StDQHn13o8KfvzRR8jPzcV/PvigQy2EU1JTsbS4GB+8/z6SU1LCH7darTj55NG46abrD/hnRxIGACKiHu7BBx/BZ5+tMaxnl1VpuOgUG1wzHdDckfHkvytdDy0JjD/HhmUF8Thj1N6PCmZkZGDnzp0ocLmwdMkSpKalwWaz7fZr9yQ7JwevvPgiXn7hBUOIkCQJOTk5KCjI79B7iiQMAEREPdjbb7+LpUufNmz6a2jSceRBElyzHIBdQCDC97EpFTqOODzUQvi2yXbsqNThD+w+BMTHx8Nut+PRhx7CosJCNDU1GZ7i9yY1LQ0fr16NJ4uLkZKaGj4BIAgCNE1DYWE+UlKSO+ttmY4BgIioh9q4cSPmzcuFqrZOnauqjqCqwzXLjkOPkaA0RObT/66UBh2iACy8Mw6Pzo6DrgPu+t0vCUiShOycHCx/7TXk5+bih3XrkJWdvdefb3c4sH3bNixdsgRer9ew+c9isWDevDkYPvz4Tn1PZmMAICLqoZzOPKiqBl1vLZQ7KjXMn+7AGefIUCqjo/i38AdDjYOmTLFj2cJ4DB4oYXvFnvcuZOfk4Id165Cfm4vlr72GrOzs5k2QRoIgwGqx4MniYqz/4QckNV9EBIQ2/V100QWYOPGyLnlPZmIAICLqgXRdx6+/bjKcU99WpmHmhTKum2qHP8I2/e0PpUzD6LGho4KTzmo+KriHfQEpqaloamrCosJCPPrQQ7A7HIiLizN8TVZWFp5ZuhTvvvUWsrKywh+3WKwYMuRo5OXN69L3YxYGACKiHkbXAZtNNqz7V9Vq+PNwS2jdv0lHdD37/5FSrSM7S8TiBfFYeI0D1bU6mny7f1eyLCM1LQ1LlyxBgcuF0p07kZ4eOsefmZWFt958E88+/bRhmUAURTgcjh616W9XDABERD2KDkkSDcXfp+jISBHhmuVAcrYIRTFxeJ1I8QKBRh233RSHJ3LjkZokoKJm70cF/7NqFfJzc/HxRx/hiEGDsG7tWjxZXAy7LBuWBwRBQGHhAhx00IDuejvdLnpbGBER0S5a2vi29q/XdR3uBh1FNzpw/CgrlG7u89/VtOajguefbwvfI7DiswD65+y+LXBmZia2btmChbm52LF9O779+muU7tyJtOYZASC07n/11TNw+umndtfbMAVnAIiIeghdDxUvTTM2+5n9VzsmXCr3uOLfllKh45jBoaOCN0ywY1u5hkBw90sCCYmJsFiteLK4GN+tXfuH4n/qqeNw7bWzumvopmEAICLqAXRdh80mIxhsbfazs1LDZf9nw9yZdgSro33Vf9+UOh12WcA9zjg8dJsDSgCoa9z9+7ZYLHDssiFQkiT0798XBQULu2vIpmIAICKKcrquw2q1Gdb96xp1DD28udmPCLTpANyjKf7QBsHpVzmwLD8eRx4k7rGFcFuCIEBVVRQULEB8fNw+v74nYAAgIopiP/20AQCgqq3H/YJqqGmOa6YD/Y+QoHjMGp15lDINfzrNitfvS0DfbBHSLtWupctfC0myID8/D8OGDe3GUZqLAYCIKEr5/X44nbmwWKyG/villTpyZ9pxypk2KFU9f+p/Txp3apCkUPOjtm0CJMkCXdchSaF98FarDZdddikuueQik0ZqDgYAIqIoNXduHjZv3mJ4+i8p03DtpTJm/tUecTf8dTdJBJISBchWAT5/awJITEzAW28tx2GHHYq4uAQcf/wwOJ1/M3Gk5mAAICKKQsXFT2LlylWGdf9Kt4YzRlrgmmUH9rD5LZYIAoD+EgYdJMLS5tC7z+fHQQcNwLJlT+Gcc87EggV5po3RTAwARERRZvXqj/Dww48ain+TT0efzFCzn7hUEYp/Lz8gRmg6gGoN/XLEcGcEURTh9/vwww/rkZCQgAULctG/fz8zh2kaBgAioiiyY8cOOJ25QJtmvrquo7FJh2uWA0NGWKDU8ekfAOLsApRqHV4FaPSGPqZpGgRBgNfrM3dwEYABgIgoijideaivb4Smta7vl5RpuPNKOy64yAalnMU/LEHAvH968fYnASTGheYARFHEwIGHYvTokYaNk7GIAYCIKEoUFt6NtWvXGZr97KjQMPlsGXfMcCAQZdf7diU5R0TRMz7881UF/XPa9vgXUVKyHcFg8A9HAWMN7wIgIooCr7zyGl566TVD8a9t0DD8aAmumXZA06HF9qb/MDlDwH9WKMhf4kPvzNYiHyr4Op577ilYLCx/nAEgIopw69Z9j7y8fMNxv0BQh2wT4JrlQK9DJChNJg4wgsgJArb8oiKvOLTGb5FaA4AkWVBYmI+jjhpk1vAiCgMAEVEE83g8cDpzIUmSYc26vFpH3kw7xp5qgxIDff7bwyIBCOpwFXnx4yYVSfGtxd9qtWHKlEk477y/mDfACMMAQEQUwZzOPGzbtgNqm2b+W0tV3Hi5jKmT2OynLSldQEGxD6+s8qNXZmt5s1isGDnyRNx++y0mji7ycBGEotKcOQvMHgLFqLvv7r6mMY8+uhj//e9qw3n/8moN555sg2uWAzqP+4XJOSJeetGHvy/zGTb9iaKEjIx0FBTw34xdcQaAiCgCrVy5Co8/vsRQ/D1eHYf0EeGaZYc1QYA/sJcfEEPkVAHffOqHq8iL1CRhl939OgoL85GZmWHa+CIVAwARUYT5/fff4XTmQtdbp/c1TYdXCTX7OfJYC5R6Pv0DgGwHaks1uIp8qKnTYbcZ1/1nz74No0adZOIIIxcDABFRhHE68+DzKYZmP9vKNTivsuOc82UoFSz+QHOvf7sAV7EXq78NIj2ltaRZrTace+7ZmDJlsnkDjHAMAEREESQvLx8//rgBwWDrkb/t5RqmnSfj5mkO+Mu56a+FLVvEo0t9KH5DQb/stpv+LBg06Aiu++8DNwESEUWI5557EcuXv4lAoHVxv6ZOw+hhFrhmOQC/jhjvXhsmZwpY8ZaCgie96Jtt7PRnsVhY/NuBAYCIKAJ89dXXuOuuewwFXgnoSIoX4JplR0ZfEUoNqz8AyIkCNq1X4SrywmoRIImt6/6iKKKwMB+HHTbQxBFGBy4BEBGZzO12w+nMgyiKaLnlT9d1VNXqyJvlwEnjbCz+zawWQPeGmv1sLNGQEGfc9DdjxjSMH3+miSOMHgwAREQmmzs3D+Xl5YZmPyVlGm67Qsaky2U2+2lDTA1t+lv+oR856W03/Vlx8smjcdNN15s4uujCJQCKSt3ZjIWoKz344CP47LM1hnX/sioNF51iQ94sBzQ3n/xbyDkili3z4cHnjTf8SZKEnJwcFBTkmzi66MMZACIik7zzzrtYuvRpQ7OfBo+OIw+WQpv+ZAGB4F5+QAyR0wR8vjrU7CcztbXZjyAI0DQNhYX5SElJNnmU0YUBgIjIBBs3boLTmQtVbZ3eV1UdQU2Ha6Ydhx4jQWng0z8AyA6gsiTU7KfRq8NmbV33t1gsmDdvDoYPP97EEUYnBgAiom6m63q4+Lft9rejUsP86Q6ccY4MpZLFHwBEAYBVQF6RF2u+DyItydjs56KLLsDEiZeZN8AoxgBARNTNnM5c/PrrJkOzn21lGmZeKOO6qXb4uekvzJot4oEnvHj6LcVw3t9isWLIkKORlzfPxNFFNwYAIqJu9NRTz+Ltt1cY1v2rajX8eXhzs58mHXz2D5GzBPx7uYK7lvoMnf5EUYTD4eCmvw5iACAi6iaffbYGf//7A4bi71N0ZKaIcM1yIDlbhKKYOMAIIicJ2LA2CFeRF3F2AWKbZj+CIKCwcAEOOmiAiSOMfgwARETdoKKiAk5nLoDWQqbrOtwNOlyz7Dh+lBUKj/wBAGxWwN+gw1Xkw9ZSDfEOY7Ofa6+9GqeffqqJI+wZGACIiLqB05mH6mo3NM3Y7Gf2X+24dIIMhZf8hAnJAlxFXrz9iR9ZacZNf6eeOg7XXjvLxNH1HAwARERd7N5778eXX36NYLC12c/OSg2XnWHD3Jl2BKv45N9CzhGxdJkPj7ykYEAvKfxxSZLQv39fFBQsNHF0PQs7ARIRdaF//etNLFv2gqH41zXqGHq4BNdMByACbToAxzQ5XcDHq/xYUOxDTrpxzV9VVRQULEB8fJyJI+xZOANARNRFfvppQ/N5/9bjfkFVhygCrlkO9D9CguIxcYARRI4DdvymIq/YC39Ah9XSGgAkyYKFC10YNmyoiSPseRgAiIi6gN/vh9OZC0GQoLe547e0UkfuDDtOOcMGhVP/AABJBCAAriIfvt2gIiXRuO4/ceIEXHzxheYNsIdiACAi6gJz5+Zh8+Ythqf/kjIN114qY+YUO2/4a8OSKWDREh+eX6GgT5ax2c8JJxyHOXNmmzi6nosBgIiokxUXP4mVK1cZzvtXujWcOdIC1yw7wB7/YXK2gNdfU7DoaZ/hhj9RFJGcnISCggUmjq5nYwAgIupEq1d/hIcfftRQ/Jt8Ovpkhpr9xKWKUPx7+QExRE4R8P2XQbiKfEhKaL3hL0RAYWE+evfuZdr4ejoGACKiTrJjx47mZj+tT/i6rqOxSYdrlgODR1ig1PHpHwBkG+Cp0uAq8qK0SoNDNjb7ueWWGzBu3FgTR9jzMQAQEXUSpzMP9fWN0LTW9f2SMg13TrXjgotsUMpZ/MMSBLiKfVj5RRCZqcZNf2ed9X+YPn2aiYOLDQwARESdoKDgbqxdu85w3n9HhYbJZ8u4Y7oDAV7vGybniCh62od/vqoY1v0lyYKBAw/hJT/dhI2AiIg66JVXXsPLL79mKP61DRqGHy2FNv1pOjRu+gcAyBkCPlihIP8JH3pnGpv9ADoKChbAamVp6g6cASAi6oDvvluHvLx8w3G/QFCHbBPgmuVAr4MlKE0mDjCCyAkCtvyiwlXkAwBYpLbNfiQUFubjqKMGmTW8mMMAQER0gBobG+F0zockGZv9lFfryJvpwNhTbVCqOfUPABYJQFCHq8iLH39TkRRv3PQ3ZcpknHfeX8wbYAxiACAiOkBz5+Zh+/ZSqG2a+W8tVXHT5TKmTpLZ7KcNKV1AQbEPr6zyo1dG201/VowceSJuv/0WE0cXmxgAiIgOwKOPLsbq1R8ZzvuXV2s492QbXLMc0HncL0zOEfHSywr+vmzXZj8S0tPT2ezHJAwARET7aeXKVXj88SWG4u/x6jikrwjXLDssCQL8gb38gBgipwr45lM/XEVepCbt2uxHR2FhPjIzM0wbXyxjACAi2g+///47nM5c6Hrr9L6m6fAqOlwzHTjyWAuUej79A4BsB2pLNbiKfKip02G3Gdf9Z8++DaNGnWTiCGMbAwAR0X5wOvPg8ymGZj/byjU4r7LjnPNlKBUs/gAgCADsAvKKvFj9bRDpKcZmP+eeezamTJls3gCJAYCIqL3y8vLx448bEAy2HvnbXq5h2nkybp7mgL+cm/5a2LJF/GOpD0uWK+iX3faGPwsGDTqC6/4RgN0WiIja4fnnX8Ty5W8iEGhd3K+p0zB6mAWuWQ7Ar0Pnwz8AQM4UsOItBQVPetG3TfEXBBEWi4XFP0IwABAR7cNXX32NwsJ7DAVe8etIihfgmmVHRl8RSg2rPwDIiQI2rVfhKvLCZhEgia3r/qIoorAwH4cdNtDEEVILLgEQEe2F2+2G05kHURTRcsufruuoqgvd8HfSOBuLfzOrBdC8oWY/G0s0JMQZN/3NmDEN48efaeIIqS0GACKivZg7Nw/l5eWGZj8lZRpuu0LGFZez2U9bYqqABcVeLP/Qj5x0Y7Ofk08ejZtuut7E0dGuGACIiPbgwQcfwWefrTGs+5dVabjoVBvyZjmgufnk30LOEbHsBQUPPr/rDX8ScnJyeMNfBGIAICLajXfeeRdLlz5taPbT4NFx5MESXDMdgE1AILiXHxBD5DQBn38YavaTmdra7EcQBGiahsLCBUhJSTZ5lLQrBgAiol1s3LgJTmcuVLV1el9VdaiaDtcsOw49RoLSyKd/AJAdQEWJhrxiHzxeHTZr67q/xWLBvHlzMHz4CSaOkPaEAYCIqA1d18PFv223vx2VGuZNd+CMs2UolSz+ACAKACwCXEVefP5DEKlJxmY/F110ASZOvMy8AdJeMQAQEbXhdObh1183GZr9bCvTMPNCGddNtcPPTX9h1mwRDzzpxdNvKeib1bbZjxVDhx6DvLx5Jo6O9oV9ACgqzZnDRiLU+f73v4349ddfDev+VbUaThne3OynSQef/UPkLAH/Xq7grqU+Q6c/URQRF+fgpr8owBkAIiIAFRXl2LDhJ0Px9yk6MlNEuGY5kJwlQlFMHGAEkZMEbFgbhKvIizi7ALFNsx9BEFBYmI8BA/qbOEJqDwYAIop5Pp8Xa9d+A6C1kOm6DndDaNPfcaMsUGr57A8AshXw1+twFfmwtVRDvMPY7Oe6667BaaedYuIIqb0YAIgo5q1d+y2CQRWaZmz2M3uKHZdOkKGUs/iHJQtwFXvx9id+ZKUZN/2ddtqfcc01M00cHO0PBgAiimnr1/8At9uNYLC12c/OSg2XnWHD3Bl2BKtY/FvIOSKefNaHR15SMKCXFP64JEno378f1/2jDDcBElHMKinZgq1btxiKf12jjqGHS6FNfyLQpgNwTJPTBXy0yo/8JT7kpBvX/FVVRWHhAsTFOUwcIe0vzgAQUUxyu91Yu/ZbqGrrcb+gqkMUAdcsB/ofLkHxmDjACCLHATt+C93w5w/osFpaA4AkWbBwoQtDhw4xcYR0IBgAiCjmqKqK7777GpJkgd7mjt/SSh25M+w45QwbFE79AwAkEYAAuIp8+PZnFSmJxnX/iRMn4OKLLzRvgHTAGACIKOasXfstvF7F8PRfUqbhuktlzJxi5w1/bVgyBSxa4sPzKxT02aXZzwknHIc5c2abODrqCO4BoKh09915Zg+BolRR0RN4551yw3n/SreGM0da4JplBxr45N9Czhbw+msKFj3tM9zwJ4oikpOTUFDAhlzRjDMARBQzPvxwNR555DFD8W/y6eiTFWr240gVofj38gNiiJwi4Psvg3AV+ZCU0HrDX0io2U/v3r1MGx91HAMAEcWE7du3w+nMBdo089V1HY1NOlwzHRg8wgKljk//ACDbAE+VBleRF6VVGhyysdnPLbfcgHHjxpo4QuoMDABEFBPmzs1DQ4MHmta6vl9SpuHOqXZccJGNzX7aShCQV+TDyi+CyEw1bvo766z/w/Tp00wcHHUWBgAi6vEKCu7GunXfG87776jQMPlsGXdMdyDA633D5BwRjz/tw+LXFMO6vyRZMHDgIVz370G4CZCIerRXXnkNL7/8mqH41zZoGH60FNr0p+nQuOkfACBnCPhghYKFS3zonWls9gPoKChYAKvVat4AqVNxBoCIeqzvvluHvLx8w3G/QFCHbBPgmuVAr4MlKE0mDjCCyAkCtvyiwlXkAwTAIrVt9iOhsDAfRx01yMQRUmdjACCiHqmxsRFOZx4kSTI0+ymvDm36G3uqDUo1p/4BwCIBCOpwFXnx428qkuKNm/6mTJmM8877i3kDpC7BAEBEPdLcuXnYvn0H1DbN/LeWqrjpchlXTpLZ7KcNKV1AQbEPr6zyo1dG201/VowcOQK3336LiaOjrsIAQEQ9zqOPLsbq1R8ZzvuXV2s4d5wNrlkO6DzuFybniHjpZQV/X7Zrsx8J6enpvOGvB2MAIKIeZeXKVXj88SWG4u/x6jikrwjXTDssCQL8gb38gBgipwr45hM/XEVepCYZm/3ouo7CwnxkZmaYOELqSgwARNRj/P7773A6c6HrrdP7mqbDp+hwzXLgyGMtUOr59A8Ash1wl2rIK/ahpl6H3WZc9//b327DqFEnmThC6moMAETUYzidefD5FEOzn23lGuZeZcc558lQKlj8AUAQANgFuIq8+OjbINKTjc1+zj33bEyZMtm8AVK3YAAgoh4hLy8fP/64AcFg65G/7eUapp0n4+ZpDvjLuemvhS1bxD+W+rBkuYJ+2W1v+LNg0KAj2OwnRrAREBFFveeffxHLl7+JQKB1cb+mTsOYYRa4ZjkAvw6dD/8AADlTwIq3FBQ86UXfNsVfEERYLBYW/xjCAEBEUe2rr75GYeE9hgKv+HUkJQhwzbIjo68IpYbVHwDkRAGb1qtwFXlhswiQxNZ1f1EUUViYj8MOG2jiCKk7cQmAiKKW2+2G05kHURTRcsufruuorgs1+zlxnI3Fv5nVAmjeULOfjSUaEuKMm/5mzrwK48efaeIIqbsxABBR1Jo7Nw/l5eWGZj8lZRpuvULGFRPZ7KctMVXAgmIvln/oR066sdnPuHGjceON15k4OjIDAwARRaUHH3wEn322xrDuX1ql4aJTbcib5YDGJ/8wOUfEs88rePD5XW/4k9CrVy8UFCw0cXRkFgYAIoo677zzLpYufdrQ7KfBo2PQwVJo059NQCC4lx8QQ+Q0AWs+9GNBsReZqa3NfgRBgKZpKCxcgOTkJJNHSWZgACCiqLJx4yY4nblQ1dbpfVXVoWo6XLPsOPRoCUojn/4BQHYAFSUaXEU+eLw6bNbWdX+LxYL58+fihBOON3GEZCYGACKKGrquh4t/225/Oyo1zJvuwBlny1AqWfwBQBQAWELNfj5fH0RqkrHZz8UXX4jLL59g3gDJdAwARBQ1nM48/PrrJkOzn21lGmZdKOO6qXb4uekvzJot4oEnvXj6LQV9s9o2+7Fi6NBjkJvrNHF0FAkYAIgoKjz11DN4++13Dev+VbUaThne3OynSQef/UPkLAH/Xq7grqU+Q6c/URQRF+fgDX8EgAGAiKLAp59+hr///UFD8fcpOjJTRbhmOZCUJUJRTBxgBJGTBGxYG4SryIs4uwCxTbMfQRBQWJiPAQP6mzhCihQMAEQU0crLy+F05gIwXlVb26DDNdOO40ZZoNTy2R8AZCvgr9fhKvJha6mGeIex2c91112N0047xcQRUiRhACCiiDZ3bh5qamqhacZmP3dMsePSCTKUchb/sGQBecVevPOJH1lpxk1/p532Z1xzzSwTB0eRhgGAiCLWokV/x9dff4NgsLXZz85KDZedYcPcGXYEq1j8W8g5Ip581od/vKSgfy8p/HFJsqB//35c96c/4GVARBSRli//F5577kVD8a9r1DH08OZmPyLQpgNwTJPTBXy0yo8FxT7kpBvX/FU1iMLCBYiLc5g4QopEnAEgoojz448/wenMg6q2HvcLBnWIIuCa5UD/wyUoHhMHGEHkOGDHb6Eb/gJBHVZLawCQJAsWLnRh6NAhJo6QIhUDABFFFJ/PB6dzPkRRgt7mjt/SKh15M+w45QwbFE79AwAkEYAAuIp8+PYXFSmJxnX/iRMn4OKLLzRvgBTRGACIKKK88MJL2LZtp+Hpv6RMw3WXypgxxc4b/tqwZApYtMSH51co6JNpbPbj9/sxZ85sE0dHkY4BgIgiytChQ+D1NoX/XFal4cxRFrhm2YEGPvm3kLMFvP6agkVP+ww3/ImiaOiXQLQnDABEZJqWKf6W1r4eTxMmT54KSbKEP983W8QT8+PhSBOhsK4BAOQUAeu+CMJV5ENSQusNfyECAAGGDxHtBgMAEZnmX//6N4DQzXQAEB8fhxEjhoen/wVBwI+/qiir1tBUw6l/AJBtQGOVBlexF6VVGhyysdmPqgZZ/KldeAyQiEyxYsX7WLToARQVLUFDgwcjRhyHM888A99+uzb8NaIIHNRPRJVbx6BDBPh8+i5PuzEoQYDrQS9WfRE0TP1brTb4/Qp/P9RuDABEZIohQwbD6/Wgrs4NQRDwwQersXLlfwxfo2mAV9Hx8xYVY0+zQYjxDYByjojHi71Y/JpiKP6SZIHfz8sQaP8wAFBUmjNngdlDoA764os14f+taRp0PYDa+iYkJdhRW9+EtJQEAECjFygp1YBGHboeuzMAcoaAD1YoWLjEh96ZxmY/Lc2SYvV3QweGewCIqFu1bPxTVRWKooSLek1tI0YeNxAPuSZh6FH9sbPcDQDolSGgtEqD0qjDnh6b/2TJCQK2/KLCVeQDBMAitW32E2r7y+JP+ys2/zYRkWlCT6xBeDweSJIULlypyfG4aPxwnHX5mSi6exom/OUkbNtZDVEAXnzfj3NvacTnawKQc2Lrny2LBCCow1XkxY+/qUiKN276mzJlMos/HZDY+ptERKbTdQ0+nw9+vx+a1rqm/+uvJfh1cyk8m7egV1YKFt91JebdeB5q6hqRmRLEuo1BTJ7vwbJlPsiZAqwxsoAppQsoKPbhlVV+9Mpou+nPipEjR+D2228xcXQUzRgAiKhbCYKIQCBguN5X13WcMm4o8udPhK7r8HgV1Dd4cedtF+GxhVOQlpIANdAIXQeuu6cJrkVNEARATuzZT75yjogXX1Lw92d3bfYjIT09nTf8UYcwABBRt9H10Ea+X3/9Obx2Hfo48N9PfsJ3X/2KJl+o24+m66gurcGlF43BkkVXYdxJg1BeWYPemQLufdaHyfM9+N8mFXJmzwwBcqqArz/xw1XsRWqysdmPrusoLMxHZmaGiSOkaMcAQETdRhBChSwpKRlqm7t8G5t8uMd5GY7981BIovGfpepyN44bfBCKF03DNZNPxc7yWvTOBN75NIBJ8z14b4Ufco6InhQDZDvgLtXgKvLBXa/DbjOu+9955+0YNeokE0dIPQEDABF1qy1bfsfmzf8zrP/7fAHsKHOjeuMOw8dbuOs8sFktuC/vCiyaMwGBQBAOmw/bKzRcmefBo4u9sKUIkOXufCddQxAA2AW4irz4aG0Q6cnGZj/nnXcO/vrXSeYNkHoMBgAi6ha6rkHXdZSUbDVc86uqGmw2C0YePxDpGUkQxd3/s+QPBFFdVY9rZp6Fonum4fCDc9DYWIekBAF/e7gJN+d7UFuvQ06N7rkAW7aIfyz1YclyBf2y297wZ8GgQUdg4UKXeYOjHiVG9tFST3P33XlmD4EOUH5+IV5++dXwnz1eBXk3n4/Lrzwd1ZvL9vn91aU1OPPUYTikfxbufuwtvPL2l+jXKw1P/MuP/23T4JrlwAmjrVDKo69roJwpYMVbCgqe9KJvm+IvCCKsVhsKC7npjzoPZwCIqNusX/8j3nzzbcM0f119E7aXuVH1Wyk0rX3X/VZX1WNA3wwU3T0Vd157DiqrG5CRHMRXPwYxaV4jXn7JBzlLCJ2hjxJyooBN64NwFXlhswiQxNaZDFEUUVi4AAMHHmriCKmnYQAgom4zePAxiItzhHe0a5qOrIxk/OW0Y5GRlghRbP/0faPHhyavgvmzL8Ej+ZOREG+HrnugBIBZdzWh8H4vJCsgx0f+koDVAmheHa4iHzaWaEiIM276mznzKpx11hkmjpB6IgYAIuo2t956BxobPeE9AFaLBH8giKyMJJRX1Rn2BrSHqumoLnPjisvGYcmiaRh53EBUVbuRky6g8EkvpuZ6sLVEhZwR2SFATBXgKvZi+Yd+5KQbm/2MGzcaN954nYmjo56KAYCIus2gQYPg83nDf3bXeXDMEX0R55Ah26wH3NK2usyNkccdhiWLrsJVl43DjjI3+mYJeOO/AUya58F/PvBHbAthOUfEs88reOj5XW/4k9CrVy8UFCw0cXTUk0Xm3wgi6nHq6urx8MP/gCi2LswnJthx1GG9kd0vE4GAupfv3rfq2kYkxNvxUP5kFM6+BB6vgkSHgs3bNfw1twlFS7yQ0wTIto6+k84jpwlY86EfC4q9yExtbfYjCAI0TUNh4QIkJyeZPErqqRgAiKjLaZqG5OQkJCcbj/ltLqlEIKBCbfTu1/r/nviUAKqrG3Djtefg8bunYkDfDDR56xFnB259wIs7Cprg8eqQk81fEpAdQEWJBleRFx6vDpu1dUwWiwXz58/FCSccb+IIqadjACCiLieKItas+QJ1dQ3hu+slUcTgI/si75YLoPiDnfp61aU1+MuZJ6B40VU459RhKC2vQb9sAf98VcGkeR788H0QcrZ5IUAUAFhCzX4+X68iNcnY7Ofiiy/E5ZdPMG18FBsYAIioy+3cWYoZM66Grrce/wsEVaz97jd8/f3vCKodm/7fnerKOhx+cA6K7pmKW2eehfKqOmSnqfjkuyAmzffgjdcUyNkC9tB3qEtZs0U88KQXT7+loG9W22Y/Vgwdegxyc53dPyiKOQwARNTlevfuhfj4BMP6/9YdVbjPNQlnnT28w+v/e1Lf6EUgoCJ/zgQ8mDcJVqsFktCE+kYd0wuacO9DXlgdgBzXJS+/W3KWgH8vV3DXUp+h058oioiLc/CGP+o2DABE1OVeffV1+Hze1ul/ScSQQf1w/ZWno67c3aWvHVQ1VJfX4srJp2LJomk4bvBBqKl1Iys1NAU/M68JO0s1yOldvyQgJwn46dtQs584u2DY9yAIAgoL8zFgQP8uHwcRwABARF2o5Vx/dXU1VFUN73JX/AF8/d1v+Oq734BuusevurQGJ580CEsWXYXJF43B9tIa9M0W8NL7fkya58HHHwW69KigbAWUeh2uYh+2lmqIdxib/Vx33dU47bRTuuz1iXbFAEBEXaal4FdXV0MQxHAg0DVg9rXn4OTzT0Yw2DXT/7tTXdOAtJQEPLZwChbcdiHqG7xITvDj599VTJ7vwdKnvJAzBNisXfDiSaEZh3c/8SMrzbjp77TT/oxrrpnVBS9KtGcMAETUZQKBABobG/HCCy9DVVt3+pdX1WHdhhKUr9/U7v7/ncXr86OmrhG33Xg+Ft91JXplpsCvNMBqAW66zwvnXU0IBkPT9Z1FzhHx5DIf/vGygv69WvdBSJIF/fv347o/mYIBgIi6jNVqRUJCAvr06YO2U/29s1PxYO4VyO6T3inn//eXroeWBC74y0kovvcqnD52MMoqatAnS8CDL/gwab4HP28IQs7q+NjkdAEfrfJjQbEPOenGNX9VDaKwcAHi4hwdfh2i/cUAQERd6plnlqGqqgZA65P+xo0l2FxSgaoy9373/+9M1RW1OOaIviheNA03TP0/lFbUoleGjlVfBjBpvgdv/TvUQvgAOxRDjgN2/KbCVeRFIKjDamn9QZJkwcKFLgwdOqST3g3R/mEAIKIuNXjwMfD5msJ/9vr8uOjcURh/yjBYrZYD7v/fWerqQ2O7e/5E/H3eRACALHlR6dYxbYEHD/6jCbZEAfJ+PqRLIgABcBX58O0vKlISjev+EydOwMUXX9hZb4NovzEAEFGX2bmzFLfeeoeh/a9ss6C8qg71jV74fH4TR9cqEFRRXVGLGdP+D8X3TMPRR/RFXX0t0pIEOB/z4jqXB5VVOuS09ocVS6aARUt8eH6Fgj6ZxmY/J5xwHObMmd0Vb4Wo3RgAiKjLvPHGctTVNUBt0+lve6kbw4cegpT0JFgs0l6+u/tVl9bg1JMHY8miqzDhLydhR1kN+mWLePbt0FHBz9e076ignC3gtVcVLHraZ7jhTxRFJCcnoaBgQVe+DaJ2YQAgoi5z1FGDoCi+8J89TQou/cuJeKBgCmqr600c2Z5VV9ejV1YKFt91JebdeB5qahuRlhTAuo1BTJ7vwXPLfJAzBVgtrd8TPt6o65BTBKz7IghXsQ9JCcIuSxyhZj+9e/fq3jdFtBsMAETU6XRdh9vtxg033ApBaLv2LaG0vBa1VfVQ/AETR7h3Hq+C+gYv7rztIjxWMAVpKQlQA43QdeDae5rgurcJggDIiQJUNXSTn5whwJ4gQPfocBV5UValwSEbm/3ccssNGDdurInvjKiVZd9fQkS0fwRBQGpqKoYOPQbfffd9+OPV7kZcP2UwUjKSUF1ea94A20HTdVSX1uDSi8bg4P6ZuOext/D+R+vRt1ca7n3Gh/9t0/Bkbjwqa3U0eDT8ulXF1lINW0s1vLcmgIP7tC5vWK02nHXW/2H69GkmviMiIwYAIuoSr7++HL/8sin8Z0kUkJzowHGDD0JdVT10XTf9BEB7VJe7cfzgg1G86Crc89hbKHr+Q/TOTMK/Pwpg1NR6rN2gYkAfER6vjnqPjn7ZoqH4S5IFAwcewnV/ijgMAETUJY444nAoig+iKEJVVfgDGhqbFHiafJBlK4Kqtu8fEiHcdR7YrBbcl3cFDu6fhUX/fAsJDhUerwNHHiJBEgGHLCAjpfV7RFGArrde8mO1dkV/YaIDxz0ARNQlFi8ugiRZwicAHA4bzj39OIyf+Oe2PYGihj8QRHVVPa6deRaK7pmGww7OQVllHQCgbZYRRTE8s2G3O1BYuACDBh1pxpCJ9oozAETUqVRVhSRJKCsrRzAYaO4BoEPXdbz+7lewT3sQs68+GwcflIXq6gazh7vfqktrcOapw3BI/yzc/dhbeOXtL9GvVxoEQQj3O7DbHSgocCE+Pp6b/ihiMQAQUaeSJAmKosDn8wEQoGmt1wAf1C8Tr73zFX7fVonZV5+NU04Ziuoyt7kDPgDVVfUY0DcDj989FUMH9cO9j7+LlKQ4CIKAxYv/AZvNhhNPHGFqm2OifeESABF1uiuumIJt27ZD1zUIggBJav2npld2CjZs2oFrnU/jyWc+QHpmMqwR1hBob3Rdb97ACMgp8bj8vJFIT4mHrus49dQ/Y+zYMTjxxBEAEBWbHCl2MQAQUadpWe/3er3h6381Tf/Dk3BKUhwgAHfe8zLm3f0ydADJiXHdPdz90lL4kxPjkNEnA1+t24zrbi3GzDuXYtvOGgDA6tWfYP58FwBA06JnkyPFJi4BEFGnkSQJ5eXlqK6uga6HiqYkGWcAWtisFuRkJuOfz34QXhIYOvQQVFfUdv/A2yElKQ6W5Hiseu8bPL98DUoravHxF7+gT04qemWnAAD8fgWVlRUAYLj/gCgSMQAQUafRdR1bt5agqckLQcAfpv93JQgC+vZKw38++ykcAs4/9yTUlLsRKcvnwaCK7N5peOn1z/DGim9QXlWHteu3IDszGQf1y/zD13/++Ve4996/Y/bs26Om1wHFJkZUIuo0giAgEAhA01Romh7+b1+y0pNQXlmHG3KfxYOPvYXU5Hg47LZuGPGeBYMq0rNS4FMCqK1qwCNL38d/Pv0JO8vc6NsrzbBvIXT0L/TPqSCIsNnk5v/N4k+RizMARNQpVFWFpmkoLFwEoHX632IR2/U0Hx8nw6FZsfDhf+H3kkrccfXZ6Nc7DdXuxi4euVEwqCK7TwY2/28HdmyrxLhLCtGvdzp+2rgDvbNTDV/bMs0futVQwMMP348TTjgeCQkJ3TpmogPBAEBEnUKSJEiShCFDjsGWLVsAwDAD0J6nYVEU0bdXGl54cw1+3xYKASeffEy3HBUMqiqy+2Zi088laNi4HaPOz8ehA7JQ1+BFYp3HUPxbmv3Y7TIUJYDi4sUYOnQI7HZ7l4+TqLMwABBRp/nqq6+xatWH4V3/kiQe0DR47+xUrNuwFdc6n8bsq8/GXy//E+qq6zu9fXDLGn16rzT88uMWlFfWY/QF+Rg0sDc0XUejx4eczGTD60qShISEBDQ2NuKpp57AkUceAavVyjP/FHUYAIio04wYMRxWqxTug9/e6f/dSU2Oh+IP4PbCF/H7tkrcPms84hwy6hu9HR5nS+EXBAHpB2Xjrdc+xYUzH8aQQf0QHyfD6/MjMy3xD4FDEERomoZXX30Rffr0hiiK4eN+XO+naMMAQESd5vbb58Dr9UFtLpz7M/2/O7LNiuyMJDy09L3wksDRR/VHdXMP/v3VUvhFQUDa4X3wwRtr0PT1Rty/ZAWyMpKg+INITY43FP5QUBChaSqGDz8OI0YMR79+fcOFn8f9KFoxABBRp/jww4/w4Ycfwe9XAIQ2xnXGU7EgCOjXKx0rPvw+HALOOXsEqktr2v0zwoVfFJB2ZD+8/fyHaPh4Pebd9xqCQQ3+QBApSXHh4NLyuqIoQZJEjBw5AqNHj8bkyVeEP8/CT9GOAYCIOsWYMaMgCKHWvy1d8wCh087CZ2cmo2RnNa6f/wx+31aJG6b+HzweH3xKYJ/fK4oi0o7sj9eXvofa975F3v2vw2KREAiqSE50IE5vPXIY6l1gQVycA6NHj8To0aNw0UUXhD/Ps/3UUzAAEFGnmDTpSgSDoaOALU/bQOeujSfG26GqGvIeeB1bmmcDsjOSUFPr+cPXtmzKy+ifieee/Q/cyz/DggeXIynRAUEQEGe3QQfCexRC4xWRnZ2FUaNOwujRo3DWWWf84eey+FNPwQBARJ3iyCOPwA8/rA93/+uqQilJIvrmpOHpVz8O9Qu45myMPGkQqstDRwVbCn9yYhysaQn4aNV3uG7eM8hIS0R8nAzZZoVss6Jlb2Ko8AvQNA2ybEVh4QKMHHlSl4ydKJIwABBRh5WUbMOrr74OQRChqiokSQ8/WXdVEOiTk4Yv1/2Ga+c+jTuuPhsTLz0ZtZW1iI9zwJqeiOWvfIyX//0ldpa7EeeQYbNaYLO2/pMXWsMPXVdss1lx3333wG6XWfwpZjAA0AEbNGiI2UOgiKJD00JP1N21QS49NQEer4JbFz6P37dVwumahOeWrMCrb3+F0opa/La1Amkp8aHbB1tGqeuwWCwIBoOwWm1QVRV+v4Kbbrq1W8ZMFCkYAIiow3Rdh9VqQzAYMHysO9bL7bIVss2CRYvfxoaN2/Hjxh0oq6hFUmIcsjKSDF8rSRYEg4Hw7EQg4G8eI9f1KfYwABBRp9A0Fbqud8nmv71pCRoD+qTjuw1boak6UpPjDa/f+tQfgCAAqhps/jy791Hs4kFWIuqwlqt/BUFovhin+19f03RAN4aOlsLf5isBCGjNBnzyp9jFAED79PPPP5g9BIpgoel/GcFg0PCxrn7NFpqmIxDQEAxqho9bLFYAoVsKAQHBYDCmCj//3tK+MABQu/AfE9qblrX/rp7+bynwLT9fVTWoqta8DBD6vM0mQxCAQCCA1kZEQCwd3+ffV2oP7gGgduM/KrQ7d9zxN6xYsTLcAAjoug2AbQu/pumGJ36bzYa+fXujtrYOK1euRkJCQvPXqpCk7l+WIIp0DABEdECCwSDKyyuwcuV/wxfjhDbVCZ1e/FsCRTCoQteNSwCiKEKSLEhIiMe7774V/nhL4WfxJ9o9BgAiOiAWiwV9+vSGw2GHoiiQJLFLz/9rmh6+XRAIFX6LxYqsrAyMGDEcY8aMBsDCT9ReDABEdMAKC++GovgB6BCEUPHv6PR/2+9vuVQoENDC+wtEUYLVakW/fn0xcuRJGDduTLj4A2DhJ2onBgAiOmBHHXUUFOUlQ8E/0OLfUvhD/wHBYGhnf8sGPlEUYbXacMQRh2H06FEYM2YUjjvu2M56K0QxhwGAiA5IZWUlcnMXhK//lSQBmnbgT/9tCz8Q2ugnCKEnelGUMHz4cRg7dgxOOulEDBp0ZGe+FaKYxABARAckMzMThx56CDZt2gQAB1T82x7rU9VQ4W956pckqflyoVB/AZcrF/379+vEd0AU2xgAiOiAvPzyq9i8eQtUVYPFIu1X8d91ur/lPD8AWCwSRFFsbiyk4e23/wW/38/iT9TJGACIaL+VlGzDXXfdi2AwAFEU9nv6v+10f8vO/pap/mAwAIvFgo8++gANDY049NBDmrv5EVFnYgAgov3Wv38/9OmTgy1bSprP5bdv81/Lk39LI5+Wj1mtNmhaEImJ8fjgg/dRXV2NrKwspKenA+DOfqKuwABARPvtsccWo7S0snnz377P/u96tK+li5/FYoEoihg2bBieffZJlJeXIy7OAZstBwALP1FXYgAgov3Wr18/+HxN4f77u7tcZ9d+AJqmh9f5W5r49OqVg9raWsybdycAhJ/4jTf4EVFX4N8yItpvr776GiTJAk0L/mEG4I+FX4Oqhqb7W87yDxjQH01NTXjvvbcMX8/CT9R9+LeNiNotGAzCYrE0r/2Hdv/vjiAAmhYKA6qqQxQF2GwyBg48FKqq4o03Xgl/raZpnOonMgEDABG1m8ViQUVFBfx+f/N0vgBJEg3n+XW95VhfaJ1flu0YPPgoCIKAZcueDv8s9uwnMhcDABHtl3XrfoDPF7r8x2IRwicAQuf5Q+v8uq5DFEXIsh2HHXYonnvumfD3s/ATRQYGACLaL3FxDgSDQaiqClGUIIpCc9Fv6QYowm6XceyxQzB69CiMGjUSQGiqP3R1Lws/USRgACCidlFVFR6PB7fd9jdIkgSrVYSqauFmPqHCb8eYMSdhzJgxGDVqJPr16xv+/q68KpiI9h8DABG1iyRJSEpKwvPPP4PrrrsJO3fuaL6wR4TDIeP0008JX9aTkZFu9nCJaB8YAIhovwwceCjef/9tAMB5512Io48+GqeddgpOPHEE4uPjTR4dEbUXAwARHbCJEy/HJZdcxOl9oijEAEBEB2zChEvMHgIRHSDGdiIiohjEAEBERBSDGACIiIhiULcHgCFDhhiuDVu/fn13D4GIiMh0u9a/XetjV+MMABERUQxiACAiIopBDABEREQxiAGAiIgoBjEAEBERxSBTAgBPAhARUSwz+wQAwBkAIiKimMQAQEREFIMYAIiIiGKQaQGA+wCIiCgWRcL6P8AZACIiopjEAEBERBSDTA0AXAYgIqJYEinT/wBnAIiIiGJSxAUAzgIQEVFPFGn1zfQAYOb0BxERkVnMrn+mBwAiIiLqfhERALgZkIiIerJI2vzXIiICABEREXWviA0AnAUgIqKeIFLrWcQEgEiYDiEiIupqkVLvIiYA7E6kpiYiIqL2iOQ6FlEBIFJSERERUVeIpDoXUQEA4IkAIiLqGSJx539bERcAdochgIiIokk01K2IDACRlpKIiIg6IhLrWkQGgN2JhjRFREQULfUqYgPA7tJStPxSiYgoNu2uTkXi0z8QwQEAYAggIqLoEU3FH4jwAABE9i+PiIhoTyK9fkV8ANgdzgIQEVEkica6FBUBgEsBREQUqaJt6r9FVAQAgCGAiIgiT7QWfyCKAgDAEEBERJEjmos/EGUBAGAIICIi80V78QeiMAAADAFERGSenlD8gSgNAHvCEEBERF2pJ9WZqA0Ae0pbPen/HCIiihx7qi/R+PQPAFE56F398MMP+q4fGzx4sBlDISKiHqinTPu3FdWDb2t3IQBgECAiogPX057624raJYBdcUmAiIg6U08u/kAPCgAAQwAREXWOnl78gR60BNDWnpYDAC4JEBHRnu3tgbEnFX+ghwaAFtwXQERE7RULT/1t9cg31RZnA4iIaG9i6am/rR77xtraWwgAGASIiGLRvvaH9eTiD8RIAGjBIEBERLFe+FvExJvc1b6CAMAwQETUk7TnNFisFP4WMfVm22pPCAAYBIiIoll7j4HHWvEHYjgAtGhvEAAYBoiIosH+9H6JxcLfImbf+O7sTxgAGAiIiCLB/jZ7i+Wi3xZ/Cbuxv0FgVwwGRESdr6NdXVn4jfjL2IeOhgEiIjIPi/6e8RezHxgGiIgiH4t++/CX1AEMBERE5mPBPzD8pXUBBgMios7HQk9EREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREREFNX+H9FysttnAXEoAAAAAElFTkSuQmCC"
_ICON_PATH_CACHE = None


def _icon_path():
    """Write the embedded icon to a temp PNG once and return its path."""
    global _ICON_PATH_CACHE
    if _ICON_PATH_CACHE and os.path.exists(_ICON_PATH_CACHE):
        return _ICON_PATH_CACHE
    fd, path = tempfile.mkstemp(suffix=".png", prefix="wb_card_icon_")
    with os.fdopen(fd, "wb") as fh:
        fh.write(base64.b64decode(_ICON_B64))
    _ICON_PATH_CACHE = path
    return path


# ── helpers ──────────────────────────────────────────────────────────────────
def _clean(value) -> str:
    """None / NaN / 'nan' -> '' so nothing prints 'nan' or '(nan)'."""
    if value is None:
        return ""
    if isinstance(value, float) and value != value:
        return ""
    s = str(value).strip()
    if s.lower() in ("nan", "none", "nat", "<na>"):
        return ""
    return s


def _get_groups(event_type, zg_override=None):
    if isinstance(zg_override, list):
        return zg_override
    return EVENT_ZIELGRUPPE.get(event_type, [])


def _card_color(groups) -> str:
    has_a  = "A" in groups
    has_pn = ("P" in groups) or ("S" in groups)
    has_pa = "PA" in groups
    if has_a and has_pn:
        return _GREEN
    if has_a:
        return _YELLOW
    if has_pn:
        return _SALMON
    if has_pa:
        return _BLUE
    return _GREY


def _audience_text(event_type, groups) -> str:
    if event_type in _SPECIAL_AUDIENCE:
        return _SPECIAL_AUDIENCE[event_type]
    return " / ".join(lbl for code, lbl in _AUDIENCE_LABELS if code in groups)


def _rgb(hex6: str):
    from pptx.dml.color import RGBColor
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _flat(shape):
    """Kill the inherited drop shadow: empty effectLst + neutralise theme effectRef."""
    from pptx.oxml.ns import qn
    shape.shadow.inherit = False                 # writes empty <a:effectLst/>
    style = shape._element.find(qn("p:style"))
    if style is not None:
        eref = style.find(qn("a:effectRef"))
        if eref is not None:
            eref.set("idx", "0")                 # 0 = no theme effect (no shadow)


# ── slide pieces ─────────────────────────────────────────────────────────────
def _fill_bg(slide, prs, hex6):
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(hex6)
    rect.line.fill.background()
    _flat(rect)
    return rect


def _text(slide, text, left, top, width, height, size, bold, color):
    from pptx.util import Pt
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = _rgb(color)
    return tb


def _cover_slide(prs, month_name, year):
    from pptx.util import Cm
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, prs, _WHITE)
    _text(slide, "Weiter- und Fortbildungsprogramm",
          Cm(1.6), Cm(4.2), Cm(22.2), Cm(2.6), 30, True, _GREEN)
    _text(slide, f"{month_name} {year}",
          Cm(1.6), Cm(6.9), Cm(16), Cm(1.6), 24, True, _BLACK)
    _text(slide, "Universitätsklinik für Intensivmedizin",
          Cm(1.6), Cm(8.6), Cm(18), Cm(1.0), 14, False, _BLACK)
    slide.shapes.add_picture(_icon_path(), Cm(19.6), Cm(9.3), Cm(4.2), Cm(4.2))
    return slide


def _section_slide(prs, week_no, date_range, kw):
    from pptx.util import Cm
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, prs, _GREEN)
    _text(slide, f"Woche {week_no}",
          Cm(2.0), Cm(4.7), Cm(21), Cm(2.4), 40, True, _WHITE)
    _text(slide, date_range,
          Cm(2.0), Cm(7.3), Cm(21), Cm(1.4), 22, False, _WHITE)
    if kw:
        _text(slide, f"KW {kw}",
              Cm(2.0), Cm(8.7), Cm(21), Cm(1.0), 14, False, _WHITE)
    return slide


def _add_card(slide, left, top, width, height, bg_hex, line1, line2, line3):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Cm, Pt

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    try:
        card.adjustments[0] = 0.09
    except Exception:
        pass
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(bg_hex)
    card.line.fill.background()
    _flat(card)

    # icon (vertically centred on the left)
    icon_sz  = Cm(2.7)
    icon_top = top + (height - icon_sz) // 2
    slide.shapes.add_picture(_icon_path(), left + Cm(0.35), icon_top,
                             icon_sz, icon_sz)

    # three text lines to the right of the icon
    tb = slide.shapes.add_textbox(left + Cm(3.5), top, width - Cm(4.0), height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = Cm(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    col = _text_color(bg_hex)

    first = True
    for text, size, bold in ((line1, 18, True), (line2, 14, False), (line3, 14, False)):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0)
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = _rgb(col)


def _card_lines(row):
    date_val = row.get("date")
    try:
        wd = WEEKDAY_DE_FULL.get(date_val.strftime("%A"), "")
        mn = MONTH_NAMES_DE.get(date_val.month, "")
        date_s = f"{wd}, {date_val.day}. {mn} {date_val.year}"
    except Exception:
        date_s = str(date_val)
    time_s = _clean(row.get("time"))
    line1 = f"{date_s}, {time_s} Uhr" if time_s else date_s

    resp  = _clean(row.get("responsible"))
    topic = _clean(row.get("topic"))
    if resp and topic:
        line2 = f"{topic} ({resp})"
    elif resp:
        line2 = f"({resp})"
    else:
        line2 = topic

    room = _clean(row.get("room"))
    event_type = row.get("event_type")
    zg_override = row.get("zielgruppe") if isinstance(row.get("zielgruppe"), list) else None
    groups = _get_groups(event_type, zg_override)
    aud = _audience_text(event_type, groups)
    if room and aud:
        line3 = f"{room}, Zielgruppe: {aud}"
    elif room:
        line3 = room
    elif aud:
        line3 = f"Zielgruppe: {aud}"
    else:
        line3 = ""

    return line1, line2, line3, _card_color(groups)


def _week_cards(prs, week_df):
    from pptx.util import Cm
    rows = list(week_df.iterrows())
    PER = 3
    MARGIN_X = Cm(1.4)
    card_w   = prs.slide_width - 2 * MARGIN_X
    top0     = Cm(1.2)
    gap      = Cm(0.45)
    card_h   = Cm(3.7)
    for start in range(0, len(rows), PER):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_bg(slide, prs, _GREY)
        for i, (_, row) in enumerate(rows[start:start + PER]):
            top = top0 + i * (card_h + gap)
            line1, line2, line3, bg = _card_lines(row)
            _add_card(slide, MARGIN_X, top, card_w, card_h, bg, line1, line2, line3)


# ── public entry point ───────────────────────────────────────────────────────
def export_to_pptx(schedule: pd.DataFrame, month: int, year: int,
                   output_dir: str = "/tmp") -> str:
    """Build one .pptx for the given month and return its file path."""
    from pptx import Presentation
    from pptx.util import Cm

    month_name = MONTH_NAMES_DE.get(month, str(month))

    df = schedule.copy()
    df["date"] = pd.to_datetime(df["date"], errors="coerce")
    df = df[df["date"].dt.month == month].sort_values(["date", "time"]).reset_index(drop=True)

    prs = Presentation()
    prs.slide_width  = Cm(25.4)
    prs.slide_height = Cm(14.29)

    _cover_slide(prs, month_name, year)

    if not df.empty:
        df["_iso"] = df["date"].dt.isocalendar().week.astype(int)
        week_no = 0
        for iso in sorted(df["_iso"].unique()):
            wk = df[df["_iso"] == iso].reset_index(drop=True)
            week_no += 1
            d0 = wk["date"].dt.date.min()
            d1 = wk["date"].dt.date.max()
            rng = (d0.strftime("%d.%m.%Y") if d0 == d1
                   else f"{d0.strftime('%d.%m.')} – {d1.strftime('%d.%m.%Y')}")
            _section_slide(prs, week_no, rng, int(iso))
            _week_cards(prs, wk)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_bg(slide, prs, _WHITE)
        _text(slide, f"Keine Daten für {month_name} {year}",
              Cm(2), Cm(6), Cm(21), Cm(2), 24, True, _GREEN)

    fname = f"Weiterbildungsplan_Slides_{month:02d}_{year}.pptx"
    fpath = os.path.join(output_dir, fname)
    prs.save(fpath)
    return fpath
